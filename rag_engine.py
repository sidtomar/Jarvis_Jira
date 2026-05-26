"""
rag_engine.py — Document loading & FAISS vector store for Superman app docs.

Uses OpenRouter for embeddings (proxied via OpenAI-compatible embeddings endpoint).

Supports:
- Pre-loaded docs from the /docs folder
- Runtime-uploaded files (PDF, DOCX, TXT, MD) via Streamlit file uploader
"""
import os
import tempfile
from pathlib import Path
from typing import List, Optional

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from config import OPENROUTER_BASE_URL, DEFAULT_EMBEDDING_MODEL


# ── Document loaders ──────────────────────────────────────────────────────────

def _load_pdf(path: str) -> List[Document]:
    from langchain_community.document_loaders import PyPDFLoader
    loader = PyPDFLoader(path)
    return loader.load()


def _load_docx(path: str) -> List[Document]:
    from langchain_community.document_loaders import Docx2txtLoader
    loader = Docx2txtLoader(path)
    return loader.load()


def _load_txt(path: str) -> List[Document]:
    from langchain_community.document_loaders import TextLoader
    loader = TextLoader(path, encoding="utf-8")
    return loader.load()


def _load_pptx(path: str) -> List[Document]:
    from pptx import Presentation
    prs = Presentation(path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        texts.append(row_text)
        if texts:
            slides_text.append(Document(
                page_content=f"Slide {i}:\n" + "\n".join(texts),
                metadata={"slide_number": i},
            ))
    return slides_text


def _load_file(path: str) -> List[Document]:
    """Route to the correct loader based on file extension."""
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return _load_pdf(path)
    elif ext in (".docx", ".doc"):
        return _load_docx(path)
    elif ext == ".pptx":
        return _load_pptx(path)
    elif ext in (".txt", ".md", ".rst"):
        return _load_txt(path)
    else:
        # Fallback: try as plain text
        try:
            return _load_txt(path)
        except Exception:
            return []


def load_docs_from_folder(folder: str = "docs") -> List[Document]:
    """Load all supported documents from the given folder."""
    docs_path = Path(folder)
    if not docs_path.exists():
        return []

    all_docs: List[Document] = []
    supported = {".pdf", ".docx", ".doc", ".pptx", ".txt", ".md", ".rst"}

    for file_path in docs_path.iterdir():
        if file_path.suffix.lower() in supported and not file_path.name.startswith("."):
            try:
                docs = _load_file(str(file_path))
                # Tag source metadata
                for doc in docs:
                    doc.metadata["source_file"] = file_path.name
                    doc.metadata["source_type"] = "preloaded"
                all_docs.extend(docs)
            except Exception as e:
                print(f"Warning: Could not load {file_path.name}: {e}")

    return all_docs


def load_docs_from_uploads(uploaded_files) -> List[Document]:
    """Load documents from Streamlit uploaded file objects."""
    all_docs: List[Document] = []

    for uploaded_file in uploaded_files:
        ext = Path(uploaded_file.name).suffix.lower()
        # Write to a temp file so loaders can read it
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            tmp.write(uploaded_file.read())
            tmp_path = tmp.name

        try:
            docs = _load_file(tmp_path)
            for doc in docs:
                doc.metadata["source_file"] = uploaded_file.name
                doc.metadata["source_type"] = "uploaded"
            all_docs.extend(docs)
        except Exception as e:
            print(f"Warning: Could not load {uploaded_file.name}: {e}")
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    return all_docs


# ── Text splitter ─────────────────────────────────────────────────────────────

def split_documents(docs: List[Document], chunk_size: int = 1000, chunk_overlap: int = 200) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
    )
    return splitter.split_documents(docs)


# ── Embeddings (via OpenRouter) ───────────────────────────────────────────────

def get_embeddings(api_key: str, model: str = DEFAULT_EMBEDDING_MODEL):
    """
    Return OpenAI-compatible embeddings via OpenRouter.
    OpenRouter proxies embedding models like openai/text-embedding-3-small.
    """
    from langchain_openai import OpenAIEmbeddings
    return OpenAIEmbeddings(
        api_key=api_key,
        base_url=OPENROUTER_BASE_URL,
        model=model,
    )


# ── Vector store ──────────────────────────────────────────────────────────────

def build_vectorstore(docs: List[Document], embeddings):
    """Build a FAISS vector store from the given documents."""
    from langchain_community.vectorstores import FAISS

    if not docs:
        return None

    chunks = split_documents(docs)
    if not chunks:
        return None

    vectorstore = FAISS.from_documents(chunks, embeddings)
    return vectorstore


def get_retriever(vectorstore, k: int = 5):
    """Get a retriever from the vectorstore."""
    if vectorstore is None:
        return None
    return vectorstore.as_retriever(search_kwargs={"k": k})


def retrieve_context(retriever, query: str) -> str:
    """Retrieve relevant context chunks as a single formatted string."""
    if retriever is None:
        return "No documentation context available. Generate the story based on the feature description alone."

    docs = retriever.invoke(query)
    if not docs:
        return "No relevant documentation found for this query."

    parts = []
    for i, doc in enumerate(docs, 1):
        source = doc.metadata.get("source_file", "Unknown")
        parts.append(f"### Excerpt {i} (from: {source})\n{doc.page_content.strip()}")

    return "\n\n".join(parts)
